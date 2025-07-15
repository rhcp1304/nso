import os
import re
import io

from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload, MediaFileUpload
from pptx import Presentation

SCOPES = ['https://www.googleapis.com/auth/drive']
TOKEN_FILE = 'token.json'
CREDENTIALS_FILE = 'bdstorage_credentials.json'
API_SERVICE_NAME = 'drive'
API_VERSION = 'v3'


class DriveHelper:
    def __init__(self, output_stream=None, style=None):
        self.output_stream = output_stream if output_stream is not None else print
        self.style = style
        if not self.style:
            class DefaultStyle:
                def SUCCESS(self, msg): return f"SUCCESS: {msg}"

                def ERROR(self, msg): return f"ERROR: {msg}"

                def WARNING(self, msg): return f"WARNING: {msg}"

            self.style = DefaultStyle()

    def _log(self, message, style_func=None):
        if style_func:
            self.output_stream(style_func(message))
        else:
            self.output_stream(message)

    def get_authenticated_drive_service(self):
        creds = None
        token_path = os.path.join(os.getcwd(), TOKEN_FILE)
        credentials_path = os.path.join(os.getcwd(), CREDENTIALS_FILE)

        if os.path.exists(token_path):
            self._log(f"Loading credentials from {TOKEN_FILE}...")
            creds = Credentials.from_authorized_user_file(token_path, SCOPES)

        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                self._log("Refreshing expired credentials...")
                try:
                    flow = InstalledAppFlow.from_client_secrets_file(credentials_path, SCOPES)
                    creds.refresh(flow.credentials)
                except Exception as e:
                    self._log(f"Error refreshing token: {e}. Re-authenticating...", style_func=self.style.ERROR)
                    creds = None

            if not creds:
                self._log(f"No valid credentials found. Initiating authentication flow (check your browser)...")
                if not os.path.exists(credentials_path):
                    self._log(
                        f"ERROR: '{CREDENTIALS_FILE}' not found at '{credentials_path}'. Please ensure it's in your project root.",
                        style_func=self.style.ERROR)
                    return None

                flow = InstalledAppFlow.from_client_secrets_file(credentials_path, SCOPES)
                creds = flow.run_local_server(port=0)

            self._log(f"Saving new credentials to {TOKEN_FILE}...")
            with open(token_path, 'w') as token:
                token.write(creds.to_json())

        self._log("Authentication successful.", style_func=self.style.SUCCESS)
        return build(API_SERVICE_NAME, API_VERSION, credentials=creds)

    def find_pptx_in_drive_folder(self, service, folder_id: str):
        self._log(f"Searching for a PPTX file in folder ID '{folder_id}'...")
        query = f"'{folder_id}' in parents and mimeType = 'application/vnd.openxmlformats-officedocument.presentationml.presentation' and trashed = false"

        try:
            results = service.files().list(
                q=query,
                spaces='drive',
                fields='files(id, name, mimeType)').execute()
            items = results.get('files', [])

            if not items:
                self._log(
                    f"No PPTX file found in folder '{folder_id}'. Please ensure the folder ID is correct and it contains a .pptx file.",
                    style_func=self.style.WARNING)
                return None, None
            else:
                if len(items) > 1:
                    self._log(
                        f"WARNING: Found {len(items)} PPTX files in folder '{folder_id}'. Using the first one found: '{items[0]['name']}'.",
                        style_func=self.style.WARNING)

                return items[0]['id'], items[0]['name']
        except HttpError as error:
            self._log(f"An error occurred while searching for PPTX in Drive: {error}", style_func=self.style.ERROR)
            return None, None

    def download_file_from_drive(self, service, file_id: str, destination_path: str):
        try:
            request = service.files().get_media(fileId=file_id)
            fh = io.FileIO(destination_path, 'wb')
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while done is False:
                status, done = downloader.next_chunk()
            return True
        except HttpError as error:
            self._log(f"An error occurred during file download from Drive (ID: {file_id}): {error}",
                      style_func=self.style.ERROR)
            return False
        except Exception as e:
            self._log(f"An unexpected error occurred during file download (ID: {file_id}): {e}",
                      style_func=self.style.ERROR)
            return False

    def upload_file_to_drive(self, service, file_name: str, file_path: str, mime_type: str, parent_folder_id: str):
        file_metadata = {
            'name': file_name,
            'parents': [parent_folder_id]
        }
        media = MediaFileUpload(file_path, mimetype=mime_type, resumable=True)
        try:
            file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()
            return file.get('id')
        except HttpError as error:
            self._log(f"An error occurred during file upload to Drive ({file_name}): {error}",
                      style_func=self.style.ERROR)
            return None
        except Exception as e:
            self._log(f"An unexpected error occurred during file upload ({file_name}): {e}",
                      style_func=self.style.ERROR)
            return None

    def extract_all_potential_links_from_last_slide(self, pptx_file_path: str) -> list[str]:
        if not os.path.exists(pptx_file_path):
            self._log(f"Error: PPTX file not found locally at '{pptx_file_path}'", style_func=self.style.ERROR)
            return []

        if not pptx_file_path.lower().endswith('.pptx'):
            self._log(f"Error: '{pptx_file_path}' is not a .pptx file.", style_func=self.style.ERROR)
            return []

        found_urls = set()
        url_pattern = re.compile(r'https?://[^\s\]\)\}>"]+')

        try:
            prs = Presentation(pptx_file_path)
            if not prs.slides:
                self._log(f"No slides found in '{pptx_file_path}'.")
                return []

            last_slide = prs.slides[-1]
            self._log(f"Analyzing the last slide (Slide {len(prs.slides)}) for all potential links...")

            def find_urls_in_text_content(text_frame_obj):
                for paragraph in text_frame_obj.paragraphs:
                    full_text = "".join([run.text for run in paragraph.runs])
                    for match in url_pattern.finditer(full_text):
                        found_urls.add(match.group(0).strip())

            for shape in last_slide.shapes:
                if hasattr(shape, 'action') and shape.action.hyperlink:
                    url = shape.action.hyperlink.address
                    if url:
                        found_urls.add(url)

                if shape.has_text_frame:
                    find_urls_in_text_content(shape.text_frame)
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.hyperlink.address:
                                url = run.hyperlink.address
                                if url:
                                    found_urls.add(url)

                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                find_urls_in_text_content(cell.text_frame)
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        if run.hyperlink.address:
                                            url = run.hyperlink.address
                                            if url:
                                                found_urls.add(url)

                if hasattr(shape, 'image') and hasattr(shape.image, 'hyperlink') and shape.image.hyperlink.address:
                    url = shape.image.hyperlink.address
                    if url:
                        found_urls.add(url)

            return list(found_urls)

        except Exception as e:
            self._log(f"An error occurred while extracting links from PPTX file: {e}", style_func=self.style.ERROR)
            return []