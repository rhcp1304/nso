import os
import re
import pickle
from google.auth.transport.requests import Request
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload
import io
from pptx import Presentation

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
CREDENTIALS_FILE = os.path.join(BASE_DIR,
                                'bdstorage_credentials.json')
TOKEN_FILE_PATH = os.path.join(BASE_DIR, 'token.pickle')


def authenticate_google_drive():
    SCOPES = ['https://www.googleapis.com/auth/drive']
    creds = None
    if os.path.exists(TOKEN_FILE_PATH):
        try:
            with open(TOKEN_FILE_PATH, 'rb') as token:
                creds = pickle.load(token)
            print("Loaded Drive API credentials from token file.")
        except Exception as e:
            print(f"Could not load Drive API token: {e}. Will re-authenticate.")
            creds = None

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            print("Drive API credentials expired, refreshing...")
            creds.refresh(Request())
        else:
            print(f"Initiating new Drive API authentication flow using {CREDENTIALS_FILE}...")
            if not os.path.exists(CREDENTIALS_FILE):
                print(f"Drive credentials file not found: {CREDENTIALS_FILE}")
                raise FileNotFoundError(
                    f"Drive credentials file not found at {CREDENTIALS_FILE}. Please ensure it's there.")

            try:
                flow = InstalledAppFlow.from_client_secrets_file(
                    CREDENTIALS_FILE, SCOPES)
                creds = flow.run_local_server(port=0)
            except Exception as e:
                print(f"Error during OAuth flow: {e}")
                return None
        try:
            with open(TOKEN_FILE_PATH, 'wb') as token:
                pickle.dump(creds, token)
            print(f"Drive API credentials saved to {TOKEN_FILE_PATH}.")
        except Exception as e:
            print(f"Failed to save Drive API token: {e}")

    try:
        return build('drive', 'v3', credentials=creds)
    except Exception as e:
        print(f"Error building Drive service: {e}")
        return None


def download_file_from_drive(service, file_id, destination_path):
    try:
        request = service.files().get_media(fileId=file_id)
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
            if status:
                print(f"Download progress: {int(status.progress() * 100)}%")

        fh.seek(0)
        with open(destination_path, 'wb') as f:
            f.write(fh.read())

        print(f"File '{file_id}' downloaded to '{destination_path}'")
        return True
    except HttpError as error:
        print(f"An HTTP error occurred during download: {error}")
        return False
    except Exception as e:
        print(f"An unexpected error occurred during download: {e}")
        return False


def get_market_and_zone_name_from_ppt(ppt_path):
    market_name = None
    zone_name = None
    try:
        prs = Presentation(ppt_path)
        if not prs.slides:
            print("PPT has no slides.")
            return None, None

        first_slide = prs.slides[0]
        slide_text = ""
        for shape in first_slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text += run.text + " "

        market_match = re.search(
            r"Market\s*Name\s*-\s*(.*?)(?:\s*ZONE|\s*Address|\s*STORE SIZE|$)",
            slide_text,
            re.IGNORECASE
        )
        if market_match:
            market_name = market_match.group(1).strip()

            market_name = re.sub(r'\s*\[Image \d+\]\s*', '', market_name).strip()
        else:
            print("Could not find 'Market Name - ' on the first slide.")

        zone_match = re.search(
            r"ZONE\s*:\s*(.*?)(?:\s*STATE|\s*CITY|\s*PIN CODE|$)",
            slide_text,
            re.IGNORECASE
        )
        if zone_match:
            zone_name = zone_match.group(1).strip()
            zone_name = re.sub(r'\s*\[Image \d+\]\s*', '', zone_name).strip()
        else:
            print("Could not find 'ZONE : ' on the first slide.")

        if market_name is None or zone_name is None:
            print(f"DEBUG: Extracted slide text:\n---START---\n{slide_text}\n---END---")

        return market_name, zone_name
    except Exception as e:
        print(f"An error occurred while reading the PPT: {e}")
        return None, None


def create_drive_folder(service, folder_name, parent_folder_id):
    file_metadata = {
        'name': folder_name,
        'mimeType': 'application/vnd.google-apps.folder',
        'parents': [parent_folder_id]
    }
    try:
        file = service.files().create(body=file_metadata, fields='id').execute()
        print(f"Folder '{folder_name}' created with ID: {file.get('id')}")
        return file.get('id')
    except HttpError as error:
        print(f"An HTTP error occurred during folder creation: {error}")
        return None
    except Exception as e:
        print(f"An unexpected error occurred during folder creation: {e}")
        return None


def find_or_create_folder(service, folder_name, parent_folder_id):
    try:
        query = (
            f"name = '{folder_name}' and "
            f"mimeType = 'application/vnd.google-apps.folder' and "
            f"'{parent_folder_id}' in parents and "
            "trashed = false"
        )
        results = service.files().list(
            q=query,
            spaces='drive',
            fields='files(id, name)'
        ).execute()

        items = results.get('files', [])
        if items:
            print(f"Found existing folder '{folder_name}' with ID: {items[0]['id']}")
            return items[0]['id']
        else:
            print(f"Folder '{folder_name}' not found, creating it...")
            return create_drive_folder(service, folder_name, parent_folder_id)
    except HttpError as error:
        print(f"An HTTP error occurred while finding/creating folder '{folder_name}': {error}")
        return None
    except Exception as e:
        print(f"An unexpected error occurred while finding/creating folder '{folder_name}': {e}")
        return None


def move_file_to_folder(service, file_id, target_folder_id):
    print(f"Attempting to move file ID: {file_id} to target folder ID: {target_folder_id}")
    try:
        file_metadata = service.files().get(fileId=file_id, fields='parents').execute()
        parents_list = file_metadata.get('parents', [])
        previous_parents = ",".join(parents_list)
        print(f"File '{file_id}' current parents: {previous_parents}")

        # Move the file
        updated_file = service.files().update(
            fileId=file_id,
            addParents=target_folder_id,
            removeParents=previous_parents,
            fields='id, parents, name'
        ).execute()
        print(f"Successfully moved file '{updated_file.get('name')}' (ID: {file_id}) to folder '{target_folder_id}'.")
        print(f"New parents of file '{file_id}': {updated_file.get('parents')}")
        return True
    except HttpError as error:
        print(f"An HTTP error occurred while moving file '{file_id}': {error.resp.status} - {error.content}")
        return False
    except Exception as e:
        print(f"An unexpected error occurred while moving file '{file_id}': {e}")
        return False


def process_ppt_and_create_folder(ppt_file_id, parent_folder_id):
    TEMP_PPT_PATH = 'temp_presentation.pptx'
    drive_service = authenticate_google_drive()
    if not drive_service:
        print("Authentication failed. Cannot proceed with PPT processing.")
        return

    if not download_file_from_drive(drive_service, ppt_file_id, TEMP_PPT_PATH):
        print("Failed to download PPT file. Cannot proceed with PPT processing.")
        return

    market_name, zone_name = get_market_and_zone_name_from_ppt(TEMP_PPT_PATH)

    if os.path.exists(TEMP_PPT_PATH):
        os.remove(TEMP_PPT_PATH)
        print(f"Removed temporary file: {TEMP_PPT_PATH}")

    if not market_name:
        print("Could not extract market name. Folder not created.")
        return

    target_parent_for_market = parent_folder_id
    if zone_name:
        print(f"Extracted Market Name: {market_name}")
        print(f"Extracted Zone Name: {zone_name}")
        zone_folder_id = find_or_create_folder(drive_service, zone_name, parent_folder_id)
        if zone_folder_id:
            target_parent_for_market = zone_folder_id
        else:
            print(
                f"Failed to find or create Zone folder '{zone_name}'. Market folder will be created directly under parent.")
    else:
        print("Could not extract zone name. Creating market folder directly under parent.")

    market_folder_id = create_drive_folder(drive_service, market_name, target_parent_for_market)

    if market_folder_id:
        move_file_to_folder(drive_service, ppt_file_id, market_folder_id)
    else:
        print(f"Failed to create Market folder '{market_name}'. PPT file not moved.")

