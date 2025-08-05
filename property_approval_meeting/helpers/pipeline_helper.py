import os
import re
import io
import pickle
import logging
import mimetypes

from google.oauth2.credentials import Credentials
from google.auth.transport.requests import Request
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload, MediaFileUpload
from pptx import Presentation

logger = logging.getLogger(__name__)

# --- Common Constants ---
TEMP_DOWNLOAD_DIRECTORY = 'temp_pipeline_downloads' # Standard temporary directory name for all downloads

# --- Google Drive API Constants ---
DRIVE_SCOPES = ['https://www.googleapis.com/auth/drive']
DRIVE_TOKEN_FILE = 'drive_token.json' # Specific token file for Drive
DRIVE_CREDENTIALS_FILE = 'bdstorage_credentials.json' # Your Drive API credentials file
DRIVE_API_SERVICE_NAME = 'drive'
DRIVE_API_VERSION = 'v3'

# --- YouTube API Constants ---
YOUTUBE_SCOPES = [
    'https://www.googleapis.com/auth/youtube.upload',
    'https://www.googleapis.com/auth/youtube.readonly'
]
YOUTUBE_CREDENTIALS_FILE = 'youtube_credentials.json' # Your YouTube API credentials file
YOUTUBE_TOKEN_FILE = 'youtube_token.pickle' # Specific token file for YouTube

# Base directory for credential/token files. Assuming they are in the same directory as this helper file.
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
YOUTUBE_CREDENTIALS_PATH = os.path.join(BASE_DIR, YOUTUBE_CREDENTIALS_FILE)
YOUTUBE_TOKEN_PATH = os.path.join(BASE_DIR, YOUTUBE_TOKEN_FILE)
DRIVE_TOKEN_PATH = os.path.join(BASE_DIR, DRIVE_TOKEN_FILE)
DRIVE_CREDENTIALS_PATH = os.path.join(BASE_DIR, DRIVE_CREDENTIALS_FILE)

# --- Video MIME Types and Extension Mapping (for identifying video files) ---
FILES_TO_PROCESS_MIME_TYPES = [
    'video/mp4', 'video/webm', 'video/quicktime', 'video/x-flv',
    'video/x-msvideo', 'video/x-ms-wmv', 'application/vnd.google-apps.video'
]

def get_file_extension(mime_type: str) -> str:
    """Returns the common file extension for a given MIME type."""
    extension = mimetypes.guess_extension(mime_type)
    if extension:
        return extension
    if mime_type == 'application/vnd.google-apps.video':
        return '.mp4' # Google Drive often converts to mp4 on download
    return ''

class IntegratedPipelineHelper:
    def __init__(self, output_stream=None, style=None):
        self.output_stream = output_stream if output_stream is not None else print
        self.style = style
        if not self.style:
            class DefaultStyle:
                def SUCCESS(self, msg): return f"SUCCESS: {msg}"
                def ERROR(self, msg): return f"ERROR: {msg}"
                def WARNING(self, msg): return f"WARNING: {msg}"
            self.style = DefaultStyle()

    def _log(self, message, style_func=None):
        """Internal logging helper that uses Django's stdout.write if available."""
        if style_func:
            self.output_stream(style_func(message))
        else:
            self.output_stream(message)

    # --- Authentication Methods (Shared) ---
    def get_authenticated_drive_service(self):
        creds = None
        if os.path.exists(DRIVE_TOKEN_PATH):
            self._log(f"Loading Drive credentials from {DRIVE_TOKEN_PATH}...")
            creds = Credentials.from_authorized_user_file(DRIVE_TOKEN_PATH, DRIVE_SCOPES)

        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                self._log("Refreshing expired Drive credentials...")
                try:
                    creds.refresh(Request())
                except Exception as e:
                    self._log(f"Error refreshing Drive token: {e}. Re-authenticating...", style_func=self.style.ERROR)
                    creds = None

            if not creds:
                self._log(f"No valid Drive credentials found. Initiating authentication flow (check your browser)...")
                if not os.path.exists(DRIVE_CREDENTIALS_PATH):
                    self._log(
                        f"ERROR: '{DRIVE_CREDENTIALS_FILE}' not found at '{DRIVE_CREDENTIALS_PATH}'. "
                        "Please ensure it's in your project root or helper directory.",
                        style_func=self.style.ERROR)
                    return None

                flow = InstalledAppFlow.from_client_secrets_file(DRIVE_CREDENTIALS_PATH, DRIVE_SCOPES)
                creds = flow.run_local_server(port=0)

            self._log(f"Saving new Drive credentials to {DRIVE_TOKEN_PATH}...")
            with open(DRIVE_TOKEN_PATH, 'w') as token:
                token.write(creds.to_json())

        self._log("Drive authentication successful.", style_func=self.style.SUCCESS)
        return build(DRIVE_API_SERVICE_NAME, DRIVE_API_VERSION, credentials=creds)

    def get_youtube_service(self):
        creds = None
        if os.path.exists(YOUTUBE_TOKEN_PATH):
            try:
                with open(YOUTUBE_TOKEN_PATH, 'rb') as token:
                    creds = pickle.load(token)
                self._log("Loaded YouTube API credentials from token file.")
            except Exception as e:
                self._log(f"Could not load YouTube API token: {e}. Will re-authenticate.", style_func=self.style.WARNING)
                creds = None

        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                self._log("YouTube API credentials expired, refreshing...")
                creds.refresh(Request())
            else:
                self._log(f"Initiating new YouTube API authentication flow using {YOUTUBE_CREDENTIALS_PATH}...")
                if not os.path.exists(YOUTUBE_CREDENTIALS_PATH):
                    self._log(f"YouTube credentials file not found: {YOUTUBE_CREDENTIALS_PATH}. Please ensure it's there.", style_func=self.style.ERROR)
                    raise FileNotFoundError(f"YouTube credentials file not found at {YOUTUBE_CREDENTIALS_PATH}.")

                flow = InstalledAppFlow.from_client_secrets_file(YOUTUBE_CREDENTIALS_PATH, YOUTUBE_SCOPES)
                creds = flow.run_local_server(port=0)

            try:
                with open(YOUTUBE_TOKEN_PATH, 'wb') as token:
                    pickle.dump(creds, token)
                self._log(f"YouTube API credentials saved to {YOUTUBE_TOKEN_PATH}.")
            except Exception as e:
                self._log(f"Failed to save YouTube API token: {e}", style_func=self.style.ERROR)

        self._log("YouTube authentication successful.", style_func=self.style.SUCCESS)
        return build('youtube', 'v3', credentials=creds)

    # --- Google Drive File Operations (Combined & Enhanced) ---
    def download_file_from_drive(self, service, file_id: str, destination_path: str):
        self._log(f"Attempting to download file ID: {file_id} to {destination_path}...")
        try:
            request = service.files().get_media(fileId=file_id)
            fh = io.FileIO(destination_path, 'wb')
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while done is False:
                status, done = downloader.next_chunk()
            self._log(f"Download successful: {destination_path}")
            return True
        except HttpError as error:
            self._log(f"An error occurred during file download from Drive (ID: {file_id}): {error}",
                      style_func=self.style.ERROR)
            return False
        except Exception as e:
            self._log(f"An unexpected error occurred during file download (ID: {file_id}): {e}",
                      style_func=self.style.ERROR)
            return False

    def upload_file_to_drive(self, service, file_name: str, file_path: str, mime_type: str, parent_folder_id: str):
        self._log(f"Attempting to upload file '{file_name}' to Drive folder '{parent_folder_id}'...")
        file_metadata = {
            'name': file_name,
            'parents': [parent_folder_id]
        }
        media = MediaFileUpload(file_path, mimetype=mime_type, resumable=True)
        try:
            file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()
            self._log(f"Upload successful: {file_name} (New Drive ID: {file.get('id')})", style_func=self.style.SUCCESS)
            return file.get('id')
        except HttpError as error:
            self._log(f"An error occurred during file upload to Drive ({file_name}): {error}",
                      style_func=self.style.ERROR)
            return None
        except Exception as e:
            self._log(f"An unexpected error occurred during file upload ({file_name}): {e}",
                      style_func=self.style.ERROR)
            return None

    def create_drive_folder(self, service, folder_name: str, parent_folder_id: str):
        self._log(f"Attempting to create folder '{folder_name}' under parent '{parent_folder_id}'...")
        file_metadata = {
            'name': folder_name,
            'mimeType': 'application/vnd.google-apps.folder',
            'parents': [parent_folder_id]
        }
        try:
            folder = service.files().create(body=file_metadata, fields='id').execute()
            self._log(f"Folder '{folder_name}' created with ID: {folder.get('id')}", style_func=self.style.SUCCESS)
            return folder.get('id')
        except HttpError as error:
            self._log(f"An HTTP error occurred during folder creation: {error}", style_func=self.style.ERROR)
            return None
        except Exception as e:
            self._log(f"An unexpected error occurred during folder creation: {e}", style_func=self.style.ERROR)
            return None

    def find_or_create_folder(self, service, folder_name: str, parent_folder_id: str):
        self._log(f"Searching for or creating folder '{folder_name}' under parent '{parent_folder_id}'...")
        try:
            query = (
                f"name = '{folder_name}' and "
                f"mimeType = 'application/vnd.google-apps.folder' and "
                f"'{parent_folder_id}' in parents and "
                "trashed = false"
            )
            results = service.files().list(
                q=query,
                spaces='drive',
                fields='files(id, name)'
            ).execute()

            items = results.get('files', [])
            if items:
                self._log(f"Found existing folder '{folder_name}' with ID: {items[0]['id']}", style_func=self.style.WARNING)
                return items[0]['id']
            else:
                self._log(f"Folder '{folder_name}' not found, creating it...")
                return self.create_drive_folder(service, folder_name, parent_folder_id)
        except HttpError as error:
            self._log(f"An HTTP error occurred while finding/creating folder '{folder_name}': {error}", style_func=self.style.ERROR)
            return None
        except Exception as e:
            self._log(f"An unexpected error occurred while finding/creating folder '{folder_name}': {e}", style_func=self.style.ERROR)
            return None

    def move_file_to_folder(self, service, file_id: str, target_folder_id: str):
        self._log(f"Attempting to move file ID: {file_id} to target folder ID: {target_folder_id}")
        try:
            file_metadata = service.files().get(fileId=file_id, fields='parents').execute()
            parents_list = file_metadata.get('parents', [])
            previous_parents = ",".join(parents_list)

            updated_file = service.files().update(
                fileId=file_id,
                addParents=target_folder_id,
                removeParents=previous_parents,
                fields='id, parents, name'
            ).execute()
            self._log(
                f"Successfully moved file '{updated_file.get('name')}' (ID: {file_id}) to folder '{target_folder_id}'.",
                style_func=self.style.SUCCESS)
            return True
        except HttpError as error:
            self._log(f"An HTTP error occurred while moving file '{file_id}': {error.resp.status} - {error.content}",
                      style_func=self.style.ERROR)
            return False
        except Exception as e:
            self._log(f"An unexpected error occurred while moving file '{file_id}': {e}",
                      style_func=self.style.ERROR)
            return False

    # --- PPTX Parsing Methods (Combined & Enhanced) ---
    def get_market_and_zone_name_from_ppt(self, pptx_file_path: str) -> tuple[str, str]:
        self._log(f"Extracting market and zone name from PPTX: {pptx_file_path}")
        if not os.path.exists(pptx_file_path):
            self._log(f"Error: PPTX file not found locally at '{pptx_file_path}'", style_func=self.style.ERROR)
            return None, None

        market_name = None
        zone_name = None
        try:
            prs = Presentation(pptx_file_path)
            if not prs.slides:
                self._log("PPT has no slides.")
                return None, None

            # Check first slide for Market Name and Zone
            first_slide = prs.slides[0]
            slide_text = ""
            for shape in first_slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            slide_text += run.text + " "

            # Refined regex to be more flexible and capture until end of line or specific keywords
            market_match = re.search(r"Market\s*Name\s*-\s*(.*?)(?=\s*ZONE|\s*Address|\s*STORE SIZE|\n|$)", slide_text, re.IGNORECASE)
            if market_match:
                market_name = market_match.group(1).strip()
                market_name = re.sub(r'\s*\[Image \d+\]\s*', '', market_name).strip() # Clean up image placeholders
                self._log(f"Extracted Market Name: '{market_name}'")
            else:
                market_name = 'Kolkata'
                self._log("Could not find 'Market Name - ' on the first slide.", style_func=self.style.WARNING)

            zone_match = re.search(r"ZONE\s*:\s*(.*?)(?=\s*STATE|\s*CITY|\s*PIN CODE|\n|$)", slide_text, re.IGNORECASE)
            if zone_match:
                zone_name = zone_match.group(1).strip()
                zone_name = re.sub(r'\s*\[Image \d+\]\s*', '', zone_name).strip() # Clean up image placeholders
                self._log(f"Extracted Zone Name: '{zone_name}'")
            else:
                self._log("Could not find 'ZONE : ' on the first slide.", style_func=self.style.WARNING)

            return market_name, zone_name
        except Exception as e:
            self._log(f"An error occurred while reading the PPT for market/zone: {e}", style_func=self.style.ERROR)
            return None, None

    def get_market_name_prefix_for_videos(self, pptx_file_path: str) -> str:
        self._log(f"Extracting market name prefix for videos from PPTX: {pptx_file_path}")
        if not os.path.exists(pptx_file_path):
            self._log(f"Error: PPTX file not found locally at '{pptx_file_path}'", style_func=self.style.ERROR)
            return ""

        try:
            prs = Presentation(pptx_file_path)
            if not prs.slides:
                self._log(f"No slides found in '{pptx_file_path}'.")
                return ""

            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            full_text = "".join([run.text for run in paragraph.runs])
                            if "Market Name -" in full_text:
                                market_name_value = full_text.split("Market Name -", 1)[1].strip()
                                # Only take the part after the last underscore if present, else use full value
                                if '_' in market_name_value:
                                    prefix = market_name_value.rsplit('_', 1)[1].strip()
                                    self._log(f"Extracted market name prefix: '{prefix}'")
                                    return prefix
                                else:
                                    self._log(f"No underscore found in market name: '{market_name_value}'. Using full value as prefix.", style_func=self.style.WARNING)
                                    return market_name_value
            self._log("Market Name field not found in the presentation for video prefixing.", style_func=self.style.WARNING)
            return ""

        except Exception as e:
            self._log(f"An error occurred while extracting market name prefix for videos: {e}", style_func=self.style.ERROR)
            return ""


    def extract_all_potential_links_from_last_slide(self, pptx_file_path: str) -> list[dict]:
        self._log(f"Extracting links from PPTX: {pptx_file_path}")
        if not os.path.exists(pptx_file_path):
            self._log(f"Error: PPTX file not found locally at '{pptx_file_path}'", style_func=self.style.ERROR)
            return []
        if not pptx_file_path.lower().endswith('.pptx'):
            self._log(f"Error: '{pptx_file_path}' is not a .pptx file.", style_func=self.style.ERROR)
            return []

        found_links_with_names = []
        url_pattern = re.compile(r'https?://[^\s\]\)\}>"]+') # Robust URL pattern

        try:
            prs = Presentation(pptx_file_path)
            if not prs.slides:
                self._log(f"No slides found in '{pptx_file_path}'.")
                return []

            last_slide = prs.slides[-1]
            self._log(f"Analyzing the last slide (Slide {len(prs.slides)}) for all potential links and associated names...")

            def get_text_from_cell(cell):
                text_content = []
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        text_content.append("".join([run.text for run in paragraph.runs]))
                return " ".join(text_content).strip()

            # Helper function to find URLs in a text frame, associating with a name if provided
            def find_urls_in_text_content(text_frame_obj, associated_name=None):
                for paragraph in text_frame_obj.paragraphs:
                    full_text = "".join([run.text for run in paragraph.runs])
                    # Find URLs directly in the text content
                    for match in url_pattern.finditer(full_text):
                        found_links_with_names.append({
                            'name': associated_name if associated_name else None,
                            'link': match.group(0).strip()
                        })
                    # Also check for explicit hyperlinks on runs
                    for run in paragraph.runs:
                        if run.hyperlink.address:
                            url = run.hyperlink.address
                            if url:
                                found_links_with_names.append({
                                    'name': associated_name if associated_name else None,
                                    'link': url
                                })

            for shape in last_slide.shapes:
                # Check for hyperlinks on shapes (e.g., buttons)
                if hasattr(shape, 'action') and shape.action.hyperlink:
                    url = shape.action.hyperlink.address
                    if url:
                        found_links_with_names.append({'name': None, 'link': url})

                # Check text frames in shapes
                if shape.has_text_frame:
                    find_urls_in_text_content(shape.text_frame)

                # Check tables for text frames and associated names
                if shape.has_table:
                    table = shape.table
                    name_col_idx = -1
                    store_name_col_idx = -1
                    if table.rows:
                        header_row = table.rows[0]
                        for i, cell in enumerate(header_row.cells):
                            cell_text = get_text_from_cell(cell).lower().strip()
                            if "name" == cell_text:
                                name_col_idx = i
                            elif "store name" == cell_text:
                                store_name_col_idx = i

                    for row_idx, row in enumerate(table.rows):
                        if row_idx == 0: continue # Skip header

                        current_row_name = None
                        if name_col_idx != -1 and name_col_idx < len(row.cells):
                            current_row_name = get_text_from_cell(row.cells[name_col_idx])
                        elif store_name_col_idx != -1 and store_name_col_idx < len(row.cells):
                            current_row_name = get_text_from_cell(row.cells[store_name_col_idx])

                        for cell in row.cells:
                            if cell.text_frame:
                                find_urls_in_text_content(cell.text_frame, associated_name=current_row_name)

                # Check images for hyperlinks
                if hasattr(shape, 'image') and hasattr(shape.image, 'hyperlink') and shape.image.hyperlink.address:
                    url = shape.image.hyperlink.address
                    if url:
                        found_links_with_names.append({'name': None, 'link': url})

            # Deduplicate links and prioritize names if multiple entries for same link
            unique_links_with_names = {}
            for item in found_links_with_names:
                link = item['link']
                name = item['name']
                if link not in unique_links_with_names:
                    unique_links_with_names[link] = {'name': name, 'link': link}
                elif name is not None and unique_links_with_names[link]['name'] is None:
                    unique_links_with_names[link]['name'] = name

            return list(unique_links_with_names.values())

        except Exception as e:
            self._log(f"An error occurred while extracting links from PPTX file: {e}", style_func=self.style.ERROR)
            logger.exception("Error extracting links from PPTX.")
            return []

    # --- YouTube Video Upload ---
    def upload_video_to_youtube(self, youtube_service, file_path, title, description="", tags=None, category_id="22"):
        self._log(f"Attempting to upload '{title}' from {file_path} to YouTube...")
        if not os.path.exists(file_path):
            self._log(f"Video file not found: {file_path}", style_func=self.style.ERROR)
            return None

        body = {
            'snippet': {
                'title': title,
                'description': description,
                'tags': tags if tags else [],
                'categoryId': category_id
            },
            'status': {
                'privacyStatus': 'unlisted', # Unlisted makes it accessible via link, but not searchable
                'selfDeclaredMadeForKids': False
            }
        }

        media_body = MediaFileUpload(file_path, chunksize=-1, resumable=True)

        try:
            request = youtube_service.videos().insert(
                part="snippet,status",
                body=body,
                media_body=media_body
            )

            response = None
            while response is None:
                status, response = request.next_chunk()
                if status:
                    self._log(f"Upload progress for '{title}': {int(status.progress() * 100)}%")

            video_id = response.get('id')
            self._log(f"Successfully uploaded video: '{title}' (Video ID: {video_id})", style_func=self.style.SUCCESS)
            return video_id

        except HttpError as e:
            error_message = f"An HTTP error {e.resp.status} occurred during YouTube upload of '{title}': {e.content.decode()}"
            self._log(error_message, style_func=self.style.ERROR)
            logger.error(error_message, exc_info=True)
            return None
        except Exception as e:
            error_message = f"An unexpected error occurred during YouTube upload of '{title}': {e}"
            self._log(error_message, style_func=self.style.ERROR)
            logger.error(error_message, exc_info=True)
            return None