import os
import re
import io
import shutil  # For cleaning up temp directory

import gradio as gr

# Google API Client Imports
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload, MediaFileUpload

# python-pptx Import
from pptx import Presentation

# --- Google Drive API Configuration (from your DriveHelper) ---
SCOPES = ['https://www.googleapis.com/auth/drive']
TOKEN_FILE = 'token.json'  # This file stores your user's authentication token
CREDENTIALS_FILE = 'bdstorage_credentials.json'  # Your client_secret.json renamed, obtained from Google Cloud Console
API_SERVICE_NAME = 'drive'
API_VERSION = 'v3'


# --- Your DriveHelper Class ---
class DriveHelper:
    def __init__(self, output_stream=None, style=None):
        # Use io.StringIO for output capture if not provided, otherwise use the given stream (e.g., print)
        self.output_stream = output_stream if output_stream is not None else print
        self.style = style
        if not self.style:
            class DefaultStyle:
                # Basic styling for direct print or plain text output
                def SUCCESS(self, msg): return f"SUCCESS: {msg}"

                def ERROR(self, msg): return f"ERROR: {msg}"

                def WARNING(self, msg): return f"WARNING: {msg}"

                def __call__(self, msg, **kwargs): return msg  # Fallback for unstyled messages

            self.style = DefaultStyle()

    def _log(self, message, style_func=None):
        if style_func:
            self.output_stream.write(style_func(message) + "\n")  # Add newline for write method
        else:
            self.output_stream.write(message + "\n")  # Add newline for write method

    def get_authenticated_drive_service(self):
        creds = None
        # Ensure paths are absolute or relative to where script is run
        token_path = os.path.join(os.getcwd(), TOKEN_FILE)
        credentials_path = os.path.join(os.getcwd(), CREDENTIALS_FILE)

        if os.path.exists(token_path):
            self._log(f"Loading credentials from {TOKEN_FILE}...")
            creds = Credentials.from_authorized_user_file(token_path, SCOPES)

        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                self._log("Refreshing expired credentials...")
                try:
                    flow = InstalledAppFlow.from_client_secrets_file(credentials_path, SCOPES)
                    creds.refresh(flow.credentials)
                except Exception as e:
                    self._log(f"Error refreshing token: {e}. Re-authenticating...", style_func=self.style.ERROR)
                    creds = None

            if not creds:
                self._log(f"No valid credentials found. Initiating authentication flow (check your browser)...")
                if not os.path.exists(credentials_path):
                    self._log(
                        f"ERROR: '{CREDENTIALS_FILE}' not found at '{credentials_path}'. Please ensure it's in the same directory as this script.",
                        style_func=self.style.ERROR)
                    return None

                try:
                    flow = InstalledAppFlow.from_client_secrets_file(credentials_path, SCOPES)
                    # For Gradio running in a web environment, run_local_server might not open browser.
                    # Consider using run_console() for headless or a specific port for local access.
                    # For simplicity, we keep run_local_server for typical local Gradio use.
                    self._log(
                        "Please open the URL displayed in your console to authenticate, if a browser doesn't open automatically.",
                        style_func=self.style.WARNING)
                    creds = flow.run_local_server(port=0)
                except Exception as e:
                    self._log(f"Authentication failed: {e}", style_func=self.style.ERROR)
                    return None

            self._log(f"Saving new credentials to {TOKEN_FILE}...")
            try:
                with open(token_path, 'w') as token:
                    token.write(creds.to_json())
            except Exception as e:
                self._log(f"ERROR: Could not save token.json: {e}. Check permissions.", style_func=self.style.ERROR)
                # This is critical, if token can't be saved, re-auth will be needed every time.
                return None

        self._log("Authentication successful.", style_func=self.style.SUCCESS)
        return build(API_SERVICE_NAME, API_VERSION, credentials=creds)

    def find_pptx_in_drive_folder(self, service, folder_id: str):
        self._log(f"Searching for a PPTX file in folder ID '{folder_id}'...")
        query = f"'{folder_id}' in parents and mimeType = 'application/vnd.openxmlformats-officedocument.presentationml.presentation' and trashed = false"

        try:
            results = service.files().list(
                q=query,
                spaces='drive',
                fields='files(id, name, mimeType)').execute()
            items = results.get('files', [])

            if not items:
                self._log(
                    f"No PPTX file found in folder '{folder_id}'. Please ensure the folder ID is correct and it contains a .pptx file.",
                    style_func=self.style.WARNING)
                return None, None
            else:
                if len(items) > 1:
                    self._log(
                        f"WARNING: Found {len(items)} PPTX files in folder '{folder_id}'. Using the first one found: '{items[0]['name']}'.",
                        style_func=self.style.WARNING)

                return items[0]['id'], items[0]['name']
        except HttpError as error:
            self._log(f"An error occurred while searching for PPTX in Drive: {error}", style_func=self.style.ERROR)
            return None, None

    def download_file_from_drive(self, service, file_id: str, destination_path: str):
        try:
            request = service.files().get_media(fileId=file_id)
            fh = io.FileIO(destination_path, 'wb')
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while done is False:
                status, done = downloader.next_chunk()
            return True
        except HttpError as error:
            self._log(f"An error occurred during file download from Drive (ID: {file_id}): {error}",
                      style_func=self.style.ERROR)
            return False
        except Exception as e:
            self._log(f"An unexpected error occurred during file download (ID: {file_id}): {e}",
                      style_func=self.style.ERROR)
            return False

    def upload_file_to_drive(self, service, file_name: str, file_path: str, mime_type: str, parent_folder_id: str):
        file_metadata = {
            'name': file_name,
            'parents': [parent_folder_id]
        }
        media = MediaFileUpload(file_path, mimetype=mime_type, resumable=True)
        try:
            file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()
            return file.get('id')
        except HttpError as error:
            self._log(f"An error occurred during file upload to Drive ({file_name}): {error}",
                      style_func=self.style.ERROR)
            return None
        except Exception as e:
            self._log(f"An unexpected error occurred during file upload ({file_name}): {e}",
                      style_func=self.style.ERROR)
            return None

    def extract_all_potential_links_from_last_slide(self, pptx_file_path: str) -> list[str]:
        if not os.path.exists(pptx_file_path):
            self._log(f"Error: PPTX file not found locally at '{pptx_file_path}'", style_func=self.style.ERROR)
            return []

        if not pptx_file_path.lower().endswith('.pptx'):
            self._log(f"Error: '{pptx_file_path}' is not a .pptx file.", style_func=self.style.ERROR)
            return []

        found_urls = set()
        # Updated regex to handle broader URL patterns robustly, including those with query params or fragments
        url_pattern = re.compile(r'https?://[^\s\]\)\}>"]+')

        try:
            prs = Presentation(pptx_file_path)
            if not prs.slides:
                self._log(f"No slides found in '{pptx_file_path}'.")
                return []

            last_slide = prs.slides[-1]
            self._log(f"Analyzing the last slide (Slide {len(prs.slides)}) for all potential links...")

            def find_urls_in_text_content(text_frame_obj):
                for paragraph in text_frame_obj.paragraphs:
                    full_text = "".join([run.text for run in paragraph.runs])
                    for match in url_pattern.finditer(full_text):
                        found_urls.add(match.group(0).strip())

            for shape in last_slide.shapes:
                if hasattr(shape, 'action') and shape.action.hyperlink:
                    url = shape.action.hyperlink.address
                    if url:
                        found_urls.add(url)

                if shape.has_text_frame:
                    find_urls_in_text_content(shape.text_frame)
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.hyperlink.address:
                                url = run.hyperlink.address
                                if url:
                                    found_urls.add(url)

                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                find_urls_in_text_content(cell.text_frame)
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        if run.hyperlink.address:
                                            url = run.hyperlink.address
                                            if url:
                                                found_urls.add(url)

                if hasattr(shape, 'image') and hasattr(shape.image, 'hyperlink') and shape.image.hyperlink.address:
                    url = shape.image.hyperlink.address
                    if url:
                        found_urls.add(url)

            return list(found_urls)

        except Exception as e:
            self._log(f"An error occurred while extracting links from PPTX file: {e}", style_func=self.style.ERROR)
            return []


# --- Main Processing Function for Gradio ---
def run_drive_process(google_drive_folder_id: str, temp_download_dir: str = 'temp_drive_downloads') -> str:
    """
    Encapsulates the logic to process Google Drive PPTX and videos for the Gradio UI.
    """
    output_capture = io.StringIO()  # To capture all log messages

    # Custom style class for Gradio, using HTML for colored output
    class GradioStyle:
        def SUCCESS(self, msg): return f"<span style='color: green; font-weight: bold;'>[SUCCESS] {msg}</span>"

        def ERROR(self, msg): return f"<span style='color: red; font-weight: bold;'>[ERROR] {msg}</span>"

        def WARNING(self, msg): return f"<span style='color: orange; font-weight: bold;'>[WARNING] {msg}</span>"

        def __call__(self, msg, **kwargs): return msg  # Default for plain messages

    style_instance = GradioStyle()

    # Instantiate DriveHelper, directing its output to our StringIO object
    drive_helper = DriveHelper(output_stream=output_capture, style=style_instance)

    output_capture.write(
        style_instance.SUCCESS(f"Starting process for Google Drive folder ID: {google_drive_folder_id}\n"))

    local_pptx_path = None  # Initialize to None for finally block

    try:
        # 1. Setup temporary directory
        if not os.path.exists(temp_download_dir):
            os.makedirs(temp_download_dir)
            output_capture.write(f"Created temporary download directory: '{temp_download_dir}'\n")
        else:
            output_capture.write(f"Using existing temporary download directory: '{temp_download_dir}'\n")

        # 2. Authenticate with Google Drive
        service = drive_helper.get_authenticated_drive_service()
        if not service:
            output_capture.write(style_instance.ERROR(
                "Failed to get authenticated Google Drive service. Please check console for authentication link.\n"))
            return output_capture.getvalue()

        # 3. Find PPTX in the specified folder
        pptx_drive_id, pptx_file_name = drive_helper.find_pptx_in_drive_folder(service, google_drive_folder_id)

        if not pptx_drive_id:
            output_capture.write(style_instance.ERROR(
                f"No PPTX file found in Google Drive folder '{google_drive_folder_id}'. Ensure the folder ID is correct and it contains exactly one .pptx file.\n"))
            return output_capture.getvalue()

        local_pptx_path = os.path.join(temp_download_dir, pptx_file_name)
        output_capture.write(f"Downloading PPTX '{pptx_file_name}' (ID: {pptx_drive_id}) from Drive...\n")

        # 4. Download PPTX
        if not drive_helper.download_file_from_drive(service, pptx_drive_id, local_pptx_path):
            output_capture.write(style_instance.ERROR(f"Failed to download PPTX file from Drive: {pptx_file_name}\n"))
            return output_capture.getvalue()
        output_capture.write(style_instance.SUCCESS(f"PPTX downloaded to: {local_pptx_path}\n"))

        # 5. Extract video links
        output_capture.write("Extracting potential video links from PPTX...\n")
        extracted_links = drive_helper.extract_all_potential_links_from_last_slide(local_pptx_path)
        output_capture.write(f"Found {len(extracted_links)} potential links.\n")

        google_drive_file_id_pattern = re.compile(r'drive\.google\.com/(?:file/d/|uc\?id=)([a-zA-Z0-9_-]+)')

        # 6. Process each extracted link
        for link in extracted_links:
            match = google_drive_file_id_pattern.search(link)
            if match:
                video_drive_id = match.group(1)
                output_capture.write(f"\n--- Detected Google Drive video link: {link} (ID: {video_drive_id}) ---\n")

                try:
                    file_metadata = service.files().get(fileId=video_drive_id, fields='name,mimeType').execute()
                    video_name = file_metadata.get('name', f"unknown_video_{video_drive_id}")
                    video_mime_type = file_metadata.get('mimeType', 'application/octet-stream')

                    local_video_path = os.path.join(temp_download_dir, video_name)

                    # Handle potential duplicate filenames in the local temp directory
                    counter = 1
                    original_local_video_path = local_video_path
                    while os.path.exists(local_video_path):
                        name_part, ext_part = os.path.splitext(original_local_video_path)
                        local_video_path = f"{name_part}_{counter}{ext_part}"
                        counter += 1

                    output_capture.write(f"Downloading Google Drive video '{video_name}' to {local_video_path}...\n")

                    if drive_helper.download_file_from_drive(service, video_drive_id, local_video_path):
                        output_capture.write(style_instance.SUCCESS(f"Successfully downloaded: {local_video_path}\n"))

                        # Basic check for empty or very small files (could indicate corrupted/error page)
                        if os.path.exists(local_video_path) and os.path.getsize(local_video_path) < 1024:
                            output_capture.write(style_instance.WARNING(
                                f"WARNING: Downloaded video '{local_video_path}' is very small ({os.path.getsize(local_video_path)} bytes). This might indicate an error or permission issue. Skipping upload.\n"))
                            continue

                        output_capture.write(
                            f"Uploading '{video_name}' back to Google Drive folder '{google_drive_folder_id}'...\n")
                        uploaded_file_id = drive_helper.upload_file_to_drive(service, video_name, local_video_path,
                                                                             video_mime_type,
                                                                             google_drive_folder_id)
                        if uploaded_file_id:
                            output_capture.write(style_instance.SUCCESS(
                                f"Successfully uploaded: {video_name} (New Drive ID: {uploaded_file_id})\n"))
                        else:
                            output_capture.write(
                                style_instance.ERROR(f"Failed to upload '{video_name}' to Google Drive.\n"))
                    else:
                        output_capture.write(style_instance.ERROR(f"Failed to download Google Drive video: {link}\n"))

                except HttpError as api_error:
                    output_capture.write(
                        style_instance.ERROR(f"Google Drive API error for video link {link}: {api_error}\n"))
                    output_capture.write(style_instance.ERROR(
                        "Check permissions for the video file or if the video file exists in Drive.\n"))
                except Exception as e:
                    output_capture.write(
                        style_instance.ERROR(f"An unexpected error occurred while processing video link {link}: {e}\n"))
            else:
                output_capture.write(
                    f"Skipping non-Google Drive link: {link} (This script only downloads and re-uploads Google Drive videos).\n")

    except Exception as e:
        output_capture.write(style_instance.ERROR(f"A critical error occurred during the process: {e}\n"))
    finally:
        # Clean up temporary PPTX file
        if local_pptx_path and os.path.exists(local_pptx_path):
            try:
                os.remove(local_pptx_path)
                output_capture.write(f"Cleaned up temporary PPTX file: {local_pptx_path}\n")
            except OSError as e:
                output_capture.write(
                    style_instance.WARNING(f"Could not remove temporary PPTX file '{local_pptx_path}': {e}\n"))

        # Clean up temp_download_dir itself if it was created by this run and is empty
        if os.path.exists(temp_download_dir):
            if not os.listdir(temp_download_dir):  # Check if directory is empty
                try:
                    os.rmdir(temp_download_dir)
                    output_capture.write(f"Removed empty temporary download directory: '{temp_download_dir}'\n")
                except OSError as e:
                    output_capture.write(style_instance.WARNING(
                        f"Could not remove empty temporary directory '{temp_download_dir}': {e}\n"))
            else:
                output_capture.write(style_instance.WARNING(
                    f"Temporary directory '{temp_download_dir}' is not empty after process, keeping it.\n"))

    output_capture.write(style_instance.SUCCESS("Process completed.\n"))
    return output_capture.getvalue()


# --- Gradio Interface Definition ---
iface = gr.Interface(
    fn=run_drive_process,
    inputs=[
        gr.Textbox(
            label="Google Drive Folder ID",
            placeholder="e.g., 1A2b3C4dE5f6G7h8I9j0K1l2M3n4O5p6",
            info="The ID of the Google Drive folder containing the single PPTX. Ensure your 'bdstorage_credentials.json' is in the same directory as this script."
        ),
        gr.Textbox(
            label="Temporary Download Directory (Optional)",
            value="temp_drive_downloads",
            placeholder="e.g., my_temp_files",
            info="Local directory for temporary downloads. Default: 'temp_drive_downloads'. Files will be cleaned up."
        )
    ],
    outputs=gr.HTML(  # Corrected: Changed from gr.Textbox to gr.HTML
        label="Process Output",  # This label is applied by gr.Interface
        elem_id="output_log"  # Optional ID for CSS styling if needed
        # Removed: lines=25, interactive=False, render_as="html" as they are not valid for gr.HTML
    ),
    title="Google Drive PPTX Video Re-Uploader",
    description="<p>This tool finds a PPTX in a specified Google Drive folder, extracts Google Drive video links from its last slide, downloads those videos, and then re-uploads them to the <b>same</b> Google Drive folder. <br><br><b>Important Setup:</b></p><ul><li>Ensure you have your Google Drive API <code>bdstorage_credentials.json</code> file in the same directory as this script.</li><li>The first time you run this, a browser window will open for Google account authentication. Follow the prompts and allow access. A <code>token.json</code> file will be created for future runs.</li><li>Install necessary Python libraries: <code>pip install gradio google-api-python-client google-auth-oauthlib python-pptx</code></li></ul>",
    allow_flagging="never",
    theme=gr.themes.Soft()  # Corrected theme setting for Gradio 4.x
)

if __name__ == "__main__":
    iface.launch(share=False)  # Set share=True to get a public link (temporarily)